#!/usr/bin/env python3
"""Convert slide-separated Markdown to a styled PowerPoint (.pptx)."""

from __future__ import annotations

import re
import sys
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ── 16:9 slide geometry ──────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
ML = Inches(0.65)
MR = Inches(0.65)
CONTENT_W = SLIDE_W - ML - MR
CONTENT_W_IN = 13.333 - 0.65 - 0.65  # inches, for layout math

COLOR_NAVY = RGBColor(0x1B, 0x3A, 0x5C)
COLOR_BLUE = RGBColor(0x2E, 0x75, 0xB6)
COLOR_LIGHT = RGBColor(0xF2, 0xF7, 0xFB)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK = RGBColor(0x33, 0x33, 0x33)
COLOR_MUTED = RGBColor(0x66, 0x66, 0x66)
COLOR_ACCENT = RGBColor(0xD6, 0x4A, 0x2A)


@dataclass
class TableBlock:
    rows: list[list[str]]
    bold_cells: list[list[bool]] = field(default_factory=list)


@dataclass
class SlideContent:
    title: str = ""
    subtitle: str = ""
    blocks: list[Any] = field(default_factory=list)  # str bullets/paras, TableBlock, code str


def strip_md_inline(text: str) -> str:
    text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
    text = re.sub(r"`(.+?)`", r"\1", text)
    text = re.sub(r"\$(.+?)\$", r"\1", text)
    return text.strip()


def parse_table_row(line: str) -> list[str] | None:
    line = line.strip()
    if not line.startswith("|"):
        return None
    cells = [c.strip() for c in line.strip("|").split("|")]
    return cells


def is_separator_row(cells: list[str]) -> bool:
    return all(re.fullmatch(r":?-{1,}:?", c.replace(" ", "")) for c in cells if c)


def parse_slides(md_text: str) -> list[SlideContent]:
    raw = re.split(r"\n---\n", md_text.strip())
    slides: list[SlideContent] = []
    for block in raw:
        lines = block.strip().splitlines()
        if not lines:
            continue
        sc = SlideContent()
        i = 0
        while i < len(lines):
            line = lines[i]
            if line.startswith("# ") and not sc.title:
                sc.title = line[2:].strip()
                i += 1
                continue
            if line.startswith("## ") and not sc.subtitle:
                sc.subtitle = line[3:].strip()
                i += 1
                continue
            break

        while i < len(lines):
            line = lines[i]
            if not line.strip():
                i += 1
                continue
            if line.startswith("```"):
                i += 1
                code_lines: list[str] = []
                while i < len(lines) and not lines[i].startswith("```"):
                    code_lines.append(lines[i])
                    i += 1
                if i < len(lines):
                    i += 1
                sc.blocks.append(("code", "\n".join(code_lines)))
                continue
            if line.startswith("|"):
                table_rows: list[list[str]] = []
                bold_grid: list[list[bool]] = []
                while i < len(lines) and lines[i].strip().startswith("|"):
                    cells = parse_table_row(lines[i])
                    if cells and not is_separator_row(cells):
                        table_rows.append([strip_md_inline(c) for c in cells])
                        bold_grid.append(["**" in c for c in cells])
                    i += 1
                if table_rows:
                    sc.blocks.append(TableBlock(rows=table_rows, bold_cells=bold_grid))
                continue
            if line.startswith("- "):
                sc.blocks.append(("bullet", line[2:].strip()))
                i += 1
                continue
            if re.match(r"^\d+\.\s", line):
                sc.blocks.append(("numbered", re.sub(r"^\d+\.\s*", "", line).strip()))
                i += 1
                continue
            sc.blocks.append(("para", line.strip()))
            i += 1
        slides.append(sc)
    return slides


def set_run_font(run, *, size: int, bold: bool = False, color: RGBColor | None = None, name: str | None = None):
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    if name:
        run.font.name = name


def add_rich_text(paragraph, text: str, *, size: int = 14, color: RGBColor = COLOR_DARK, base_bold: bool = False):
    parts = [p for p in re.split(r"(\*\*.+?\*\*)", text) if p]
    if not parts:
        paragraph.text = ""
        return

    if "**" not in text:
        paragraph.text = strip_md_inline(text)
        if paragraph.runs:
            set_run_font(paragraph.runs[0], size=size, bold=base_bold, color=color)
        return

    # rebuild runs for mixed bold/plain text
    while paragraph.runs:
        paragraph._p.remove(paragraph.runs[0]._r)

    for part in parts:
        run = paragraph.add_run()
        if part.startswith("**") and part.endswith("**"):
            run.text = part[2:-2]
            set_run_font(run, size=size, bold=True, color=color)
        else:
            run.text = strip_md_inline(part)
            set_run_font(run, size=size, bold=base_bold, color=color)


def fill_shape(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def add_header_bar(slide, title: str, subtitle: str = "") -> float:
    """Return y-offset (inches) where body content should start."""
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(1.05))  # rectangle
    fill_shape(bar, COLOR_NAVY)
    bar.line.fill.background()

    title_box = slide.shapes.add_textbox(ML, Inches(0.18), CONTENT_W, Inches(0.55))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    add_rich_text(p, title, size=26, color=COLOR_WHITE, base_bold=True)

    body_top = 1.15
    if subtitle:
        sub_box = slide.shapes.add_textbox(ML, Inches(0.62), CONTENT_W, Inches(0.38))
        stf = sub_box.text_frame
        stf.margin_left = stf.margin_right = 0
        sp = stf.paragraphs[0]
        add_rich_text(sp, subtitle, size=14, color=RGBColor(0xBB, 0xD4, 0xF0))
        body_top = 1.05

    # thin accent line
    accent = slide.shapes.add_shape(1, ML, Inches(body_top - 0.06), CONTENT_W, Inches(0.04))
    fill_shape(accent, COLOR_BLUE)
    accent.line.fill.background()
    return body_top + 0.12


def add_title_slide(prs: Presentation, sc: SlideContent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    fill_shape(bg, COLOR_NAVY)
    bg.line.fill.background()

    accent = slide.shapes.add_shape(1, Inches(0), Inches(3.55), SLIDE_W, Inches(0.06))
    fill_shape(accent, COLOR_ACCENT)
    accent.line.fill.background()

    # main title
    box = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.3), Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_rich_text(p, sc.title, size=36, color=COLOR_WHITE, base_bold=True)

    if sc.subtitle:
        sbox = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.3), Inches(0.6))
        stf = sbox.text_frame
        sp = stf.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        add_rich_text(sp, sc.subtitle, size=20, color=RGBColor(0xBB, 0xD4, 0xF0))

    # meta lines
    y = 4.0
    for block in sc.blocks:
        if isinstance(block, tuple) and block[0] == "para":
            mbox = slide.shapes.add_textbox(Inches(2.0), Inches(y), Inches(9.3), Inches(0.45))
            mp = mbox.text_frame.paragraphs[0]
            mp.alignment = PP_ALIGN.CENTER
            add_rich_text(mp, block[1], size=15, color=RGBColor(0xDD, 0xE8, 0xF5))
            y += 0.42


def add_closing_slide(prs: Presentation, sc: SlideContent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    fill_shape(bg, COLOR_NAVY)
    bg.line.fill.background()

    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.2))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_rich_text(p, sc.title, size=40, color=COLOR_WHITE, base_bold=True)

    if sc.subtitle:
        sbox = slide.shapes.add_textbox(Inches(1.0), Inches(3.5), Inches(11.3), Inches(0.6))
        sp = sbox.text_frame.paragraphs[0]
        sp.alignment = PP_ALIGN.CENTER
        add_rich_text(sp, sc.subtitle, size=20, color=RGBColor(0xBB, 0xD4, 0xF0))

    y = 4.5
    for block in sc.blocks:
        if isinstance(block, tuple) and block[0] == "para":
            mbox = slide.shapes.add_textbox(Inches(2.5), Inches(y), Inches(8.3), Inches(0.4))
            mp = mbox.text_frame.paragraphs[0]
            mp.alignment = PP_ALIGN.CENTER
            add_rich_text(mp, block[1], size=13, color=RGBColor(0xCC, 0xDD, 0xEE))
            y += 0.35


def estimate_block_height(block: Any, _content_width_in: float) -> float:
    if isinstance(block, TableBlock):
        n = len(block.rows)
        return 0.38 + n * 0.36
    if isinstance(block, tuple):
        kind, text = block
        if kind == "code":
            return 0.25 + text.count("\n") * 0.22 + 0.2
        lines = max(1, len(text) / 70)
        if kind == "bullet":
            return 0.32 + lines * 0.22
        if kind == "numbered":
            return 0.32 + lines * 0.22
        return 0.28 + lines * 0.22
    return 0.3


def add_table(slide, block: TableBlock, top: float) -> float:
    rows_data = block.rows
    bold_grid = block.bold_cells or [[False] * len(r) for r in rows_data]
    if not rows_data:
        return top
    nrows = len(rows_data)
    ncols = max(len(r) for r in rows_data)
    for r in rows_data:
        while len(r) < ncols:
            r.append("")
    for r in bold_grid:
        while len(r) < ncols:
            r.append(False)

    row_h = 0.36 if nrows <= 6 else 0.32
    table_h = Inches(row_h * nrows + 0.08)
    table_shape = slide.shapes.add_table(nrows, ncols, ML, Inches(top), CONTENT_W, table_h)
    table = table_shape.table

    total_w = CONTENT_W_IN
    if ncols == 5 and nrows <= 5:
        fracs = [0.10, 0.14, 0.14, 0.14, 0.48]
    elif ncols == 5:
        fracs = [0.16, 0.14, 0.14, 0.14, 0.42]
    elif ncols == 4 and nrows >= 7:
        fracs = [0.22, 0.26, 0.26, 0.26]
    elif ncols == 4:
        fracs = [0.20, 0.22, 0.22, 0.36]
    elif ncols == 3:
        fracs = [0.28, 0.36, 0.36]
    else:
        fracs = [1.0 / ncols] * ncols
    for ci, frac in enumerate(fracs[:ncols]):
        table.columns[ci].width = Inches(total_w * frac)

    highlight_rows = set()
    for ri, row in enumerate(rows_data):
        if ri == 0:
            continue
        if any(bold_grid[ri]) or any(c in {"VidCom²", "AOT", "PruneVid"} for c in row):
            highlight_rows.add(ri)

    for ri, row in enumerate(rows_data):
        row_bold = ri in highlight_rows
        for ci, cell_text in enumerate(row):
            cell = table.cell(ri, ci)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Pt(8)
            cell.margin_top = cell.margin_bottom = Pt(4)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            is_header = ri == 0
            is_first_col = ci == 0
            p.alignment = PP_ALIGN.LEFT if is_first_col or is_header else PP_ALIGN.CENTER
            cell_bold = is_header or bold_grid[ri][ci] or row_bold
            p.text = cell_text
            if p.runs:
                set_run_font(
                    p.runs[0],
                    size=11 if ncols >= 5 else 12,
                    bold=cell_bold,
                    color=COLOR_WHITE if is_header else COLOR_DARK,
                )
            if is_header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_BLUE
            elif row_bold:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFC)
            elif ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_LIGHT
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE

    return top + row_h * nrows + 0.22


def add_code_block(slide, code: str, top: float) -> float:
    lines = code.splitlines()
    h = 0.28 + len(lines) * 0.24
    box = slide.shapes.add_shape(1, ML, Inches(top), CONTENT_W, Inches(h))
    fill_shape(box, RGBColor(0xEE, 0xF2, 0xF6))
    box.line.color.rgb = RGBColor(0xCC, 0xD6, 0xE0)

    tbox = slide.shapes.add_textbox(ML + Inches(0.15), Inches(top + 0.08), CONTENT_W - Inches(0.3), Inches(h - 0.12))
    tf = tbox.text_frame
    tf.word_wrap = False
    for li, line in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Consolas"
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_DARK
        p.space_after = Pt(1)
    return top + h + 0.15


def add_text_block(slide, kind: str, text: str, top: float) -> float:
    width = CONTENT_W - (Inches(0.25) if kind == "bullet" else Inches(0))
    left = ML + (Inches(0.25) if kind == "bullet" else Inches(0))
    est_lines = max(1, len(strip_md_inline(text)) / 62)
    h = 0.28 + est_lines * 0.24
    box = slide.shapes.add_textbox(left, Inches(top), width, Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = 0
    p = tf.paragraphs[0]
    if kind == "bullet":
        p.level = 0
        add_rich_text(p, "•  " + text, size=15, color=COLOR_DARK)
    elif kind == "numbered":
        add_rich_text(p, text, size=15, color=COLOR_DARK)
    else:
        add_rich_text(p, text, size=14, color=COLOR_DARK)
    return top + h + 0.06


def render_content_slide(prs: Presentation, sc: SlideContent):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    y = add_header_bar(slide, sc.title, sc.subtitle)
    max_y = 7.15
    num_counter = 0

    i = 0
    blocks = sc.blocks
    while i < len(blocks):
        block = blocks[i]
        needed = estimate_block_height(block, CONTENT_W_IN)
        if y + needed > max_y and i > 0:
            cont_title = sc.title + "（续）"
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            y = add_header_bar(slide, cont_title, "")
            num_counter = 0
        if isinstance(block, TableBlock):
            y = add_table(slide, block, y)
        elif isinstance(block, tuple):
            kind, text = block
            if kind == "code":
                y = add_code_block(slide, text, y)
            elif kind == "numbered":
                num_counter += 1
                y = add_text_block(slide, "numbered", f"{num_counter}. {text}", y)
            else:
                y = add_text_block(slide, kind, text, y)
        i += 1


def convert(md_path: Path, pptx_path: Path) -> None:
    slides = parse_slides(md_path.read_text(encoding="utf-8"))
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for idx, sc in enumerate(slides):
        is_first = idx == 0
        is_last = idx == len(slides) - 1
        if is_first:
            add_title_slide(prs, sc)
        elif is_last and "谢谢" in sc.title:
            add_closing_slide(prs, sc)
        else:
            render_content_slide(prs, sc)

    prs.save(str(pptx_path))
    print(f"Wrote {len(prs.slides)} slides -> {pptx_path}")


def main() -> None:
    md_path = Path(sys.argv[1]) if len(sys.argv) > 1 else Path(__file__).parent / "video_token_compression_interview.md"
    pptx_path = Path(sys.argv[2]) if len(sys.argv) > 2 else md_path.with_suffix(".pptx")
    convert(md_path, pptx_path)


if __name__ == "__main__":
    main()
